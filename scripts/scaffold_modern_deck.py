# -*- coding: utf-8 -*-
"""scaffold_modern_deck.py — 可编辑 · WPS 安全 的 python-pptx 通用脚手架。

这是 python-pptx-editable-deck skill 的起点。它包含经 WPS/PowerPoint 验证的安全
原语（字体修复、实色软阴影、卡片、pill、页码、logo 嵌块等），并附一个最小示例
（封面 + Why 两页，文案均为占位）。复制进你的项目后，改 OUT / LOGO / 配色，并把
示例 build_* 换成真实内容即可。

铁律（详见 skill references/technical.md）：
  - 禁手写 a:gradFill（渐变）
  - 禁手写 a:alpha（半透明）—— 阴影用标准 a:outerShdw + 实色 srgbClr
  - 中文必须写 <a:ea>（set_run_fonts + patch_theme_fonts）
"""
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image

# ============================ 配置（按项目改这里） ============================
BASE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(BASE, "deck_output.pptx")   # TODO: 改成你的输出路径
LOGO = os.path.join(BASE, "logo.png")           # TODO: 主机/品牌 logo（透明 PNG）；无则设为 None

# ---- 颜色（明亮浅色主题，WPS 安全：纯色、无渐变/无 alpha）----
BG       = "F4F6FB"   # 浅冷灰白底
SURFACE  = "FFFFFF"   # 白色卡片
SURFACE2 = "EEF2F9"   # 次级浅面
BORDER   = "E3E8F0"   # 浅描边
TEXT     = "1E293B"   # 主文字（深板岩）
MUTED    = "64748B"   # 次要文字
FAINT    = "E6EBF3"   # 极淡装饰
BLUE     = "2563EB"   # 主强调·明亮蓝
CYAN     = "0891B2"   # 明亮青
VIOLET   = "7C3AED"   # 明亮紫
AMBER    = "EA580C"   # 活力橙
GREEN    = "059669"   # 明亮绿
SHADOW   = "C7D2E3"   # 卡片柔和阴影色（实色，WPS 安全）
ARROW    = "64748B"   # 流程箭头颜色（浅底清晰）

# ---- 字体（全 Windows 自带，≤3 种）----
FONT_BASE  = "Microsoft YaHei"   # 正文/注释/页码/标签
FONT_TITLE = "DengXian"          # 标题/卡片标题
FONT_SUB   = "DengXian Light"    # 副标题/引导语

prs = Presentation()
prs.slide_width  = int(1280 * 9525)
prs.slide_height = int(720 * 9525)
BLANK = prs.slide_layouts[6]


# ============================ 基础工具 ============================
def px(n):
    return int(n * 9525)

def Ox(name):
    return OxmlElement(name)

def set_run_fonts(run, name):
    """同时设置 latin / ea(中日韩) / cs，确保中文也用指定字体渲染。"""
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = OxmlElement('a:latin'); rPr.append(latin)
    latin.set('typeface', name)
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = OxmlElement('a:ea'); latin.addnext(ea)
    ea.set('typeface', name)
    cs = rPr.find(qn('a:cs'))
    if cs is None:
        cs = OxmlElement('a:cs'); ea.addnext(cs)
    cs.set('typeface', name)

def set_font(run, name=FONT_BASE, size=None, bold=None, italic=None, color=None, underline=False):
    set_run_fonts(run, name)
    if size is not None:   run.font.size = Pt(size)          # 用真实 point；原项目用 *0.75 自行约定
    if bold is not None:   run.font.bold = bold
    if italic is not None: run.font.italic = italic
    if underline:          run.font.underline = True
    if color is not None:  run.font.color.rgb = RGBColor.from_string(color)

def add_bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, px(1280), px(720))
    r.line.fill.background(); r.fill.solid()
    r.fill.fore_color.rgb = RGBColor.from_string(BG); r.shadow.inherit = False
    ring1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(1040), px(-120), px(360), px(360))
    ring1.line.color.rgb = RGBColor.from_string(FAINT); ring1.line.width = Pt(1.5)
    ring1.fill.background(); ring1.shadow.inherit = False
    ring2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(-120), px(560), px(320), px(320))
    ring2.line.color.rgb = RGBColor.from_string(FAINT); ring2.line.width = Pt(1.5)
    ring2.fill.background(); ring2.shadow.inherit = False

def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=4, font=None):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    base_font = font or FONT_BASE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        if isinstance(line, str):       line = [{"t": line}]
        elif isinstance(line, dict):    line = [line]
        for seg in line:
            r = p.add_run(); r.text = seg["t"]
            set_font(r, name=seg.get("font", base_font), size=seg.get("size", 14),
                     bold=seg.get("bold", False), italic=seg.get("italic", False),
                     color=seg.get("color", TEXT))
            if seg.get("underline"): r.font.underline = True
    return tb

def add_pill(slide, text, right_x, top, fill=AMBER, text_color="14181F"):
    w, h = 156, 42
    x = right_x - w
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(top), px(w), px(h))
    pill.adjustments[0] = 0.5; pill.line.fill.background()
    pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor.from_string(fill)
    pill.shadow.inherit = False
    tf = pill.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    set_font(r, size=15, bold=True, color=text_color)
    return pill

def add_logo(slide, top=28, h=92, logo_path=None):
    logo_path = logo_path or LOGO
    if not logo_path or not os.path.exists(logo_path): return
    with Image.open(logo_path) as im:
        iw, ih = im.size
    w = int(h * iw / ih)
    slide.shapes.add_picture(logo_path, px(60), px(top), px(w), px(h))

def add_header(slide, tag, tag_color=AMBER, logo_path=None):
    add_logo(slide, top=26, h=88, logo_path=logo_path)
    add_pill(slide, tag, right_x=1216, top=32, fill=tag_color, text_color="14181F")

def add_pageno(slide, n, total=6):
    add_text(slide, 1148, 686, 64, 20,
             [{"t": "%02d / %02d" % (n, total), "size": 11, "color": MUTED}],
             align=PP_ALIGN.RIGHT)

def add_ring_node(slide, cx, cy, d, color, outline=True):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx - d/2), px(cy - d/2), px(d), px(d))
    o.shadow.inherit = False
    if outline:
        o.fill.background(); o.line.color.rgb = RGBColor.from_string(color); o.line.width = Pt(1.5)
    else:
        o.fill.solid(); o.fill.fore_color.rgb = RGBColor.from_string(color); o.line.fill.background()
    return o

def add_edge(slide, x1, y1, x2, y2, color=BORDER, w=1.25):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2))
    c.line.color.rgb = RGBColor.from_string(color); c.line.width = Pt(w); c.shadow.inherit = False
    return c

def add_chip(slide, x, y, s, color):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(s), px(s))
    chip.adjustments[0] = 0.22; chip.line.fill.background()
    chip.fill.solid(); chip.fill.fore_color.rgb = RGBColor.from_string(color); chip.shadow.inherit = False
    return chip

def add_logo_chip(slide, x, y, s, color, logo_png, frac=0.6):
    add_chip(slide, x, y, s, color)
    with Image.open(logo_png) as im:
        iw, ih = im.size
    lw = int(s * frac); lx = x + (s - lw) / 2; ly = y + (s - lw) / 2
    slide.shapes.add_picture(logo_png, px(lx), px(ly), px(lw), px(lw))

def add_text_chip(slide, x, y, s, color, text):
    add_chip(slide, x, y, s, color)
    tf = slide.shapes.add_textbox(px(x), px(y), px(s), px(s)).text_frame
    tf.word_wrap = True; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    set_font(r, size=15, bold=True, color="FFFFFF")

def add_soft_shadow(shape, color=SHADOW, blur=110000, dist=38000, dir=5400000):
    """标准 a:outerShdw + 实色（无 alpha），WPS/PowerPoint 兼容。"""
    spPr = shape._element.spPr
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None: spPr.remove(existing)
    eff = OxmlElement('a:effectLst'); sh = OxmlElement('a:outerShdw')
    sh.set('blurRad', str(blur)); sh.set('dist', str(dist))
    sh.set('dir', str(dir)); sh.set('rotWithShape', '0')
    c = OxmlElement('a:srgbClr'); c.set('val', color); sh.append(c)
    eff.append(sh); spPr.append(eff)

def add_card(slide, x, y, w, h, fill=SURFACE, border=BORDER, radius=0.08):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h))
    card.adjustments[0] = radius
    card.line.color.rgb = RGBColor.from_string(border); card.line.width = Pt(1)
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor.from_string(fill)
    card.shadow.inherit = False
    add_soft_shadow(card)
    return card

def add_badge(slide, cx, cy, d, text, fill=BLUE, color="FFFFFF"):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx - d/2), px(cy - d/2), px(d), px(d))
    o.line.fill.background(); o.fill.solid(); o.fill.fore_color.rgb = RGBColor.from_string(fill)
    o.shadow.inherit = False
    add_soft_shadow(o, color="B9C4D6", blur=70000, dist=26000)
    tf = o.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    set_font(r, size=15, bold=True, color=color)


def add_click_sequence(slide, steps):
    """点击依次出现：每次鼠标点击，按顺序出现一组 shape（进入=出现/appear）。
    steps: list[list[int]]，每个内层 list = 同一次点击同时出现的 shape id。
    未列入 steps 的元素（标题/页眉）从一开始就可见。标准 <p:timing>，无渐变/无 alpha，
    WPS/PowerPoint 兼容。"""
    from lxml import etree
    sld = slide._element
    old = sld.find(qn('p:timing'))
    if old is not None: sld.remove(old)
    timing = etree.SubElement(sld, qn('p:timing'))
    tnLst = etree.SubElement(timing, qn('p:tnLst'))
    root = etree.SubElement(tnLst, qn('p:par'))
    rc = etree.SubElement(root, qn('p:cTn')); rc.set('id', '1'); rc.set('fill', 'hold')
    sc = etree.SubElement(rc, qn('p:stCondLst'))
    for dvt in ('0', 'indefinite'):
        c = etree.SubElement(sc, qn('p:cond')); c.set('srId', '0'); c.set('dvt', dvt)
    rchild = etree.SubElement(rc, qn('p:childTnLst'))
    _id = [1]
    def nid():
        _id[0] += 1; return str(_id[0])
    for group in steps:
        par = etree.SubElement(rchild, qn('p:par'))
        ctn = etree.SubElement(par, qn('p:cTn')); ctn.set('id', nid()); ctn.set('fill', 'hold')
        cs = etree.SubElement(ctn, qn('p:stCondLst'))
        c0 = etree.SubElement(cs, qn('p:cond')); c0.set('srId', '0'); c0.set('dvt', '0')
        schild = etree.SubElement(ctn, qn('p:childTnLst'))
        for spId in group:
            ppar = etree.SubElement(schild, qn('p:par'))
            pctn = etree.SubElement(ppar, qn('p:cTn')); pctn.set('id', nid()); pctn.set('fill', 'hold')
            ps = etree.SubElement(pctn, qn('p:stCondLst'))
            pc = etree.SubElement(ps, qn('p:cond')); pc.set('srId', '0'); pc.set('dvt', '0')
            pch = etree.SubElement(pctn, qn('p:childTnLst'))
            anim = etree.SubElement(pch, qn('p:animEffect'))
            anim.set('id', nid()); anim.set('filter', 'appear'); anim.set('transition', 'in')
            anim.set('presetId', '1'); anim.set('presetClass', 'entr'); anim.set('presetSubtype', '0'); anim.set('build', '0')
            cb = etree.SubElement(anim, qn('p:cBhvr'))
            cbc = etree.SubElement(cb, qn('p:cTn')); cbc.set('id', nid()); cbc.set('fill', 'hold')
            cbs = etree.SubElement(cbc, qn('p:stCondLst'))
            cbc0 = etree.SubElement(cbs, qn('p:cond')); cbc0.set('srId', '0'); cbc0.set('dvt', '0')
            tgt = etree.SubElement(cb, qn('p:tgtEl'))
            spt = etree.SubElement(tgt, qn('p:spTgt')); spt.set('spId', str(spId))
            anl = etree.SubElement(cb, qn('p:attrNameLst'))
            an = etree.SubElement(anl, qn('p:attrName')); an.text = 'style.visibility'
    etree.SubElement(timing, qn('p:bldLst'))
    etree.SubElement(timing, qn('p:extLst'))


# ============================ 示例页（占位文案，替换成真实内容） ============================
def build_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_logo(s, top=26, h=104)
    add_pill(s, "标签", right_x=1216, top=36, fill=AMBER, text_color="14181F")
    add_text(s, 74, 150, 720, 24,
             [{"t": "副标题 / 引导语", "size": 15, "color": MUTED, "font": FONT_SUB}])
    tb = s.shapes.add_textbox(px(70), px(186), px(900), px(170))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    r = p1.add_run(); r.text = "主标题"
    set_font(r, name=FONT_TITLE, size=58, bold=True, color=TEXT)
    r = p1.add_run(); r.text = " 强调"
    set_font(r, name=FONT_TITLE, size=58, bold=True, color=BLUE)
    p2 = tf.add_paragraph(); r = p2.add_run(); r.text = "一句话结论"
    set_font(r, name=FONT_TITLE, size=58, bold=True, color=TEXT)
    add_text(s, 74, 362, 820, 30,
             [{"t": "支撑主标题的那句说明文字", "size": 18, "color": MUTED, "font": FONT_SUB}])
    add_pageno(s, 1)

def build_example_page():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "SECTION · 章节", tag_color=AMBER)
    add_text(s, 90, 150, 1100, 60,
             [[{"t": "这是本页的", "size": 38, "bold": True, "color": TEXT},
               {"t": "Action Title（观点，不是话题）", "size": 38, "bold": True, "color": BLUE}]],
             align=PP_ALIGN.LEFT, font=FONT_TITLE)
    cards = [
        ("01", "卡片标题", "卡片描述文字，控制在两行内", CYAN),
        ("02", "卡片标题", "卡片描述文字，控制在两行内", VIOLET),
        ("03", "卡片标题", "卡片描述文字，控制在两行内", AMBER),
    ]
    cw, ch, gap, x0, y0 = 340, 250, 40, 90, 250
    groups = []
    for i, (no, title, desc, col) in enumerate(cards):
        x = x0 + i * (cw + gap)
        card = add_card(s, x, y0, cw, ch, fill=SURFACE, border=BORDER)
        badge = add_badge(s, x + 36, y0 + 38, 44, no, fill=col, color="14181F")
        t = add_text(s, x + 24, y0 + 78, cw - 48, 30,
                     [{"t": title, "size": 22, "bold": True, "color": TEXT}], align=PP_ALIGN.LEFT)
        d = add_text(s, x + 24, y0 + 120, cw - 48, 90,
                     [{"t": desc, "size": 15, "color": MUTED}], align=PP_ALIGN.LEFT)
        groups.append([card.shape_id, badge.shape_id, t.shape_id, d.shape_id])
    # 演示：三张卡点击依次出现（标题/页眉始终可见）
    add_click_sequence(s, groups)
    add_pageno(s, 2)


# ============================ 主题字体兜底 + 保存 ============================
def patch_theme_fonts(path):
    """把主题 major/minor 的 latin+ea+cs 也设为项目字体，防未指定文本回退宋体。
    直接回写 path（truncate 重写，不依赖 os.unlink）—— 沙箱 safe-delete 会拦截
    shutil.move 的删除回退，tmp+move 模式在受限环境会失败。"""
    import re, zipfile
    with zipfile.ZipFile(path, 'r') as zin:
        names = zin.namelist(); data = {n: zin.read(n) for n in names}
    xml = data['ppt/theme/theme1.xml'].decode('utf-8')

    def patch_block(block, latin_name, ea_name):
        new = block
        for tag, val in (('a:latin', latin_name), ('a:ea', ea_name), ('a:cs', ea_name)):
            new = re.sub(r'<%s[^>]*?/>' % tag, '<%s typeface="%s"/>' % (tag, val), new)
            new = re.sub(r'<%s[^>]*?>(.*?)</%s>' % (tag, tag),
                         r'<%s typeface="%s">\1</%s>' % (tag, val, tag), new)
        return new

    for open_tag, ln, ea in (('a:majorFont', 'DengXian', 'DengXian'),
                             ('a:minorFont', 'Microsoft YaHei', 'Microsoft YaHei')):
        m = re.search(r'<%s>.*?</%s>' % (open_tag, open_tag), xml, re.S)
        if m:
            xml = xml[:m.start()] + patch_block(m.group(0), ln, ea) + xml[m.end():]
    data['ppt/theme/theme1.xml'] = xml.encode('utf-8')
    with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for n in names: zout.writestr(n, data[n])


if __name__ == "__main__":
    build_cover()
    build_example_page()
    prs.save(OUT)
    patch_theme_fonts(OUT)
    print("saved", OUT, os.path.getsize(OUT), "bytes")
